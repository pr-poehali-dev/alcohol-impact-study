import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_color(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def set_bg(slide, color_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_color(color_hex)


def add_gradient_rect(slide, left, top, width, height, color1_hex, color2_hex=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    fill = shape.fill
    if color2_hex:
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = hex_color(color1_hex)
        fill.gradient_stops[1].color.rgb = hex_color(color2_hex)
    else:
        fill.solid()
        fill.fore_color.rgb = hex_color(color1_hex)
    return shape


def add_text_box(slide, text, left, top, width, height, font_size=18, bold=False,
                  color_hex="FFFFFF", align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_color(color_hex)
    run.font.name = font_name
    return txBox


def add_card(slide, left, top, width, height, color_hex="1E1035", border_hex="7C3AED"):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(color_hex)
    shape.line.color.rgb = hex_color(border_hex)
    shape.line.width = Pt(0.75)
    return shape


def add_bullet_list(slide, items, left, top, width, height, font_size=13, color_hex="CCCCFF"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = hex_color(color_hex)
        run.font.name = "Calibri"


def add_progress_bar(slide, label, value, left, top, width, color_hex="9333EA"):
    add_text_box(slide, label, left, top, width * 0.7, 0.22, font_size=11, color_hex="CCCCCC")
    add_text_box(slide, f"{value}%", left + width * 0.72, top, width * 0.25, 0.22,
                 font_size=11, bold=True, color_hex="C084FC")

    bar_top = top + 0.22
    bg_bar = slide.shapes.add_shape(1, Inches(left), Inches(bar_top), Inches(width), Inches(0.1))
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = hex_color("2D1B69")
    bg_bar.line.fill.background()

    fill_width = width * (value / 100)
    if fill_width > 0.05:
        fg_bar = slide.shapes.add_shape(1, Inches(left), Inches(bar_top), Inches(fill_width), Inches(0.1))
        fg_bar.fill.gradient()
        fg_bar.fill.gradient_angle = 0
        fg_bar.fill.gradient_stops[0].color.rgb = hex_color(color_hex)
        fg_bar.fill.gradient_stops[1].color.rgb = hex_color("EC4899")
        fg_bar.line.fill.background()


def handler(event: dict, context) -> dict:
    """Генерирует PPTX-файл презентации 'Влияние алкоголя на организм человека' и возвращает base64."""

    if event.get('httpMethod') == 'OPTIONS':
        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Methods': 'GET, OPTIONS',
                'Access-Control-Allow-Headers': 'Content-Type',
            },
            'body': ''
        }

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BG = "0D0020"
    BG2 = "130028"
    WHITE = "FFFFFF"
    PURPLE = "A855F7"
    PINK = "EC4899"
    ORANGE = "F97316"
    LIGHT = "E9D5FF"
    CARD = "1A0A3D"
    CARD2 = "200B45"
    BORDER = "4C1D95"
    RED = "EF4444"
    GREEN = "22C55E"
    YELLOW = "EAB308"
    CYAN = "06B6D4"
    GOLD = "F59E0B"

    blank_layout = prs.slide_layouts[6]

    # ===================== СЛАЙД 1: ТИТУЛЬНЫЙ =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 7.5, "1A0040", "0D0020")

    add_gradient_rect(sl, 0, 0, 5, 7.5, "3B0764", "0D0020")

    add_text_box(sl, "🧬", 5.9, 0.3, 1.5, 1.2, font_size=54, align=PP_ALIGN.CENTER)

    add_text_box(sl, "МУНИЦИПАЛЬНОЕ АВТОНОМНОЕ ОБЩЕОБРАЗОВАТЕЛЬНОЕ УЧРЕЖДЕНИЕ",
                 2.5, 1.3, 8.3, 0.5, font_size=9, color_hex="A78BFA", align=PP_ALIGN.CENTER)

    add_text_box(sl, "«Средняя школа №152 имени А.Д. Березина»",
                 2.2, 1.75, 8.9, 0.5, font_size=16, bold=True, color_hex=WHITE, align=PP_ALIGN.CENTER)

    line = sl.shapes.add_shape(1, Inches(5.4), Inches(2.35), Inches(2.5), Inches(0.05))
    line.fill.gradient()
    line.fill.gradient_angle = 0
    line.fill.gradient_stops[0].color.rgb = hex_color(PURPLE)
    line.fill.gradient_stops[1].color.rgb = hex_color(ORANGE)
    line.line.fill.background()

    add_text_box(sl, "Итоговый индивидуальный проект",
                 2.5, 2.5, 8.3, 0.4, font_size=13, color_hex="C4B5FD", align=PP_ALIGN.CENTER)

    add_gradient_rect(sl, 1.5, 2.95, 10.3, 1.3, "4C1D95", "831843")
    add_text_box(sl, "«Влияние алкоголя на организм человека:\nбиологические аспекты»",
                 1.7, 3.0, 9.9, 1.2, font_size=22, bold=True, color_hex=WHITE, align=PP_ALIGN.CENTER)

    add_card(sl, 3.7, 4.6, 6.0, 1.6, CARD, BORDER)
    add_text_box(sl, "Выполнил: Садомов Илья, 9Д класс",
                 3.9, 4.7, 5.6, 0.4, font_size=14, bold=True, color_hex="F9A8D4", align=PP_ALIGN.CENTER)
    add_text_box(sl, "Руководитель: Садомова Екатерина Геннадьевна",
                 3.9, 5.15, 5.6, 0.35, font_size=11, color_hex="D8B4FE", align=PP_ALIGN.CENTER)
    add_text_box(sl, "г. Красноярск, 2026",
                 3.9, 5.5, 5.6, 0.35, font_size=11, color_hex="A78BFA", align=PP_ALIGN.CENTER)

    # ===================== СЛАЙД 2: ВВЕДЕНИЕ =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 0.08, PURPLE, PINK)
    add_text_box(sl, "ВВЕДЕНИЕ", 0.4, 0.2, 3, 0.35, font_size=9, color_hex="A78BFA")
    add_text_box(sl, "Зачем это важно?", 0.4, 0.55, 8, 0.65, font_size=30, bold=True, color_hex=WHITE)

    add_card(sl, 0.4, 1.4, 5.8, 3.5, CARD, "7C3AED")
    add_text_box(sl, "⚠️  Актуальность", 0.7, 1.5, 5.2, 0.4, font_size=14, bold=True, color_hex="F9A8D4")
    add_bullet_list(sl, [
        "Алкоголь широко распространён среди молодёжи",
        "Подростки не знают реальных биологических последствий",
        "Существует миф о «безопасности» пива и коктейлей"
    ], 0.7, 1.95, 5.3, 2.0, font_size=13)

    add_card(sl, 6.5, 1.4, 6.4, 1.6, CARD, "DB2777")
    add_text_box(sl, "🎯  Цель", 6.8, 1.5, 5.8, 0.35, font_size=13, bold=True, color_hex="FB923C")
    add_text_box(sl, "Изучить механизмы действия алкоголя на организм и создать памятку для сверстников",
                 6.8, 1.9, 5.8, 0.9, font_size=12, color_hex=LIGHT)

    add_card(sl, 6.5, 3.2, 6.4, 2.5, CARD, BORDER)
    add_text_box(sl, "Задачи", 6.8, 3.3, 5.8, 0.35, font_size=13, bold=True, color_hex="C084FC")
    tasks = [
        "1. Изучить историю употребления алкоголя",
        "2. Рассмотреть биохимию этанола",
        "3. Проанализировать влияние на органы",
        "4. Провести опрос среди одноклассников",
        "5. Создать памятку «Биология трезвости»"
    ]
    add_bullet_list(sl, tasks, 6.8, 3.7, 5.8, 2.0, font_size=11, color_hex="D8B4FE")

    # ===================== СЛАЙД 3: БИОХИМИЯ =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 0.08, PURPLE, RED)
    add_text_box(sl, "БИОХИМИЯ", 0.4, 0.2, 3, 0.35, font_size=9, color_hex="A78BFA")
    add_text_box(sl, "Этанол C₂H₅OH —", 0.4, 0.55, 7, 0.6, font_size=28, bold=True, color_hex=WHITE)
    add_gradient_rect(sl, 4.45, 0.6, 4, 0.55, "BE123C", "9333EA")
    add_text_box(sl, "клеточный яд", 4.55, 0.6, 3.8, 0.55, font_size=28, bold=True, color_hex=WHITE)

    steps = [
        ("🫀", "1", "Всасывание", "20% в желудке,\n80% в кишечнике"),
        ("🩸", "2", "Кровь", "Этанол разносится\nпо всем органам"),
        ("🔬", "3", "Обезвреживание", "Происходит\nв печени"),
    ]
    for i, (icon, num, title, desc) in enumerate(steps):
        x = 0.4 + i * 4.3
        add_card(sl, x, 1.4, 3.9, 2.2, CARD2, BORDER)
        add_text_box(sl, icon, x + 0.2, 1.5, 0.6, 0.6, font_size=28)
        add_text_box(sl, num, x + 3.0, 1.45, 0.7, 0.6, font_size=32, bold=True, color_hex="3B0764")
        add_text_box(sl, title, x + 0.2, 2.15, 3.5, 0.4, font_size=16, bold=True, color_hex=WHITE)
        add_text_box(sl, desc, x + 0.2, 2.6, 3.5, 0.8, font_size=12, color_hex="C4B5FD")

    add_card(sl, 0.4, 3.85, 12.5, 2.0, "1F0A0A", RED)
    add_text_box(sl, "☠️", 0.6, 3.95, 0.7, 0.7, font_size=30)
    add_text_box(sl, "Главная опасность: Ацетальдегид",
                 1.4, 3.95, 10, 0.45, font_size=16, bold=True, color_hex="FCA5A5")
    add_text_box(sl, "В печени этанол превращается в токсичное вещество, которое:",
                 1.4, 4.42, 10, 0.35, font_size=12, color_hex=LIGHT)
    tags = ["Вызывает похмелье", "Разрушает клетки изнутри", "Накапливается при избытке"]
    for i, tag in enumerate(tags):
        tx = 1.4 + i * 3.7
        add_card(sl, tx, 4.82, 3.4, 0.55, "3B0A0A", "EF4444")
        add_text_box(sl, tag, tx + 0.1, 4.88, 3.2, 0.4, font_size=11, color_hex="FCA5A5", align=PP_ALIGN.CENTER)

    # ===================== СЛАЙД 4: МОЗГ =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 0.08, PURPLE, PINK)
    add_text_box(sl, "НЕЙРОБИОЛОГИЯ", 0.4, 0.2, 4, 0.35, font_size=9, color_hex="A78BFA")
    add_text_box(sl, "Мозг страдает первым и сильнее всего!", 0.4, 0.55, 12, 0.6, font_size=27, bold=True, color_hex=WHITE)

    add_card(sl, 0.4, 1.4, 6.0, 2.5, CARD, BORDER)
    add_text_box(sl, "Что происходит", 0.7, 1.5, 5.4, 0.4, font_size=14, bold=True, color_hex="C084FC")
    events = [
        "Этанол легко проникает в нейроны",
        "Нарушается передача сигналов между клетками",
        "Нервные клетки массово гибнут"
    ]
    for i, e in enumerate(events):
        add_card(sl, 0.7, 1.95 + i * 0.62, 5.4, 0.52, "2D1060", "7C3AED")
        add_text_box(sl, f"!  {e}", 0.9, 2.0 + i * 0.62, 5.0, 0.4, font_size=12, color_hex=LIGHT)

    add_card(sl, 0.4, 4.1, 6.0, 1.4, "1A0A00", "F97316")
    add_text_box(sl, "🧒  Подростки в зоне риска", 0.7, 4.2, 5.4, 0.4, font_size=13, bold=True, color_hex="FB923C")
    add_text_box(sl, "Мозг подростка более уязвим — разрушение происходит быстрее, чем у взрослых",
                 0.7, 4.65, 5.4, 0.7, font_size=12, color_hex=LIGHT)

    add_card(sl, 6.8, 1.4, 6.1, 4.1, CARD, "DB2777")
    add_text_box(sl, "Последствия", 7.1, 1.5, 5.5, 0.4, font_size=14, bold=True, color_hex="F9A8D4")
    consequences = [
        ("📉", "Ухудшение памяти и внимания"),
        ("🎓", "Снижение интеллекта и способности учиться"),
        ("🚶", "Нарушение координации и речи"),
        ("😶", "Деградация личности при регулярном употреблении"),
    ]
    for i, (icon, text) in enumerate(consequences):
        add_card(sl, 7.1, 2.0 + i * 0.82, 5.5, 0.7, "200B45", "4C1D95")
        add_text_box(sl, icon, 7.2, 2.05 + i * 0.82, 0.5, 0.55, font_size=22)
        add_text_box(sl, text, 7.8, 2.1 + i * 0.82, 4.6, 0.5, font_size=12, color_hex=LIGHT)

    # ===================== СЛАЙД 5: ОРГАНЫ =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 0.08, PURPLE, ORANGE)
    add_text_box(sl, "ОРГАНЫ", 0.4, 0.2, 3, 0.35, font_size=9, color_hex="A78BFA")
    add_text_box(sl, "Удар по всему организму", 0.4, 0.55, 10, 0.6, font_size=30, bold=True, color_hex=WHITE)

    organs = [
        ("🟡", "Печень", "Жировая дистрофия → гепатит → цирроз (необратимо)", "1A1000", YELLOW),
        ("❤️", "Сердце", "Склеивание эритроцитов → тромбы → инфаркты, инсульты", "1A0000", RED),
        ("🟢", "Желудок", "Ожог слизистой → гастрит, язва", "001A00", "22C55E"),
        ("🔵", "Поджелудочная", "Разрушение тканей → панкреатит", "00001A", "3B82F6"),
        ("🧬", "Репродуктив.", "Снижение тестостерона, бесплодие, поражение половых клеток", "0A001A", PURPLE),
    ]
    for i, (icon, organ, effect, bg, border) in enumerate(organs):
        x = 0.35 + i * 2.55
        add_card(sl, x, 1.45, 2.3, 3.6, bg, border)
        add_text_box(sl, icon, x + 0.7, 1.55, 0.9, 0.7, font_size=28, align=PP_ALIGN.CENTER)
        add_text_box(sl, organ, x + 0.1, 2.3, 2.1, 0.45, font_size=13, bold=True, color_hex=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, effect, x + 0.1, 2.8, 2.1, 2.1, font_size=10, color_hex="D8B4FE", align=PP_ALIGN.CENTER)

    add_card(sl, 0.4, 5.25, 12.5, 1.0, "1A1200", GOLD)
    add_text_box(sl, "🍺", 0.6, 5.35, 0.6, 0.6, font_size=26)
    add_text_box(sl, "Важно!  Пиво и коктейли — это ТОТ ЖЕ алкоголь. Одна бутылка пива = 50 г водки по воздействию на мозг.",
                 1.3, 5.4, 11.3, 0.7, font_size=13, bold=True, color_hex="FDE68A")

    # ===================== СЛАЙД 6: ОПРОС =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 0.08, PURPLE, CYAN)
    add_text_box(sl, "ИССЛЕДОВАНИЕ", 0.4, 0.2, 4, 0.35, font_size=9, color_hex="A78BFA")
    add_text_box(sl, "Результаты опроса  30 учеников  9-х классов", 0.4, 0.55, 12, 0.6, font_size=26, bold=True, color_hex=WHITE)

    surveys = [
        ("Пробовали алкоголь?", [("Да", 63, "9333EA"), ("Нет", 37, "9333EA")]),
        ("Пиво безвредно?", [("Да", 17, "9333EA"), ("Нет", 53, "9333EA"), ("Не знают", 30, "9333EA")]),
        ("Какой орган страдает первым?", [("Печень", 70, "9333EA"), ("Мозг", 13, "9333EA")]),
        ("Алкоголь влияет на умств. развитие?", [("Да", 43, "9333EA"), ("Нет / не уверены", 57, "9333EA")]),
    ]
    positions = [(0.4, 1.4), (0.4, 3.8), (6.9, 1.4), (6.9, 3.8)]
    for idx, ((title, bars), (px, py)) in enumerate(zip(surveys, positions)):
        add_card(sl, px, py, 6.0, 2.15, CARD, BORDER)
        add_text_box(sl, title, px + 0.2, py + 0.1, 5.6, 0.35, font_size=12, bold=True, color_hex="C084FC")
        for bi, (label, val, col) in enumerate(bars):
            add_progress_bar(sl, label, val, px + 0.2, py + 0.55 + bi * 0.48, 5.5, col)

    add_card(sl, 0.4, 6.1, 12.5, 0.9, "1A1400", GOLD)
    add_text_box(sl, "Вывод: Подростки знают о вреде для печени, но НЕ знают о вреде для мозга. Многие ошибочно считают пиво безопасным.",
                 0.6, 6.2, 12.1, 0.65, font_size=12, bold=True, color_hex="FDE68A")

    # ===================== СЛАЙД 7: ПАМЯТКА =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 0.08, GREEN, CYAN)
    add_text_box(sl, "ПАМЯТКА", 0.4, 0.2, 3, 0.35, font_size=9, color_hex="A78BFA")
    add_text_box(sl, "Биология трезвости: 5 главных фактов", 0.4, 0.55, 12, 0.6, font_size=28, bold=True, color_hex=WHITE)

    facts = [
        ("🧠", "1", "МОЗГ", "Алкоголь убивает нейроны. Клетки мозга не восстанавливаются. Снижается память и IQ.", "200A4A", PURPLE),
        ("🫀", "2", "ПЕЧЕНЬ", "Разрушается поэтапно: ожирение → гепатит → цирроз. Цирроз — это навсегда.", "2A1000", ORANGE),
        ("❤️", "3", "СЕРДЦЕ", "Алкоголь склеивает эритроциты. Образуются тромбы → риск инфаркта и инсульта.", "2A0010", RED),
        ("🧬", "4", "БУДУЩИЕ ДЕТИ", "Алкоголь повреждает ДНК половых клеток → риск рождения больных детей.", "0A102A", "3B82F6"),
        ("🍺", "5", "ПИВО НЕ БЕЗОПАСНО", "В пиве и коктейлях тот же этанол. Газ и сахар ускоряют всасывание яда в кровь.", "1A1200", GOLD),
    ]
    for i, (icon, num, title, desc, bg, border) in enumerate(facts):
        x = 0.3 + i * 2.55
        add_card(sl, x, 1.45, 2.3, 4.5, bg, border)
        add_text_box(sl, icon, x + 0.65, 1.55, 0.9, 0.7, font_size=28, align=PP_ALIGN.CENTER)
        add_text_box(sl, num, x + 1.7, 1.5, 0.5, 0.55, font_size=24, bold=True, color_hex="2D1B69")
        add_text_box(sl, title, x + 0.1, 2.35, 2.1, 0.55, font_size=11, bold=True, color_hex=WHITE, align=PP_ALIGN.CENTER)
        line2 = sl.shapes.add_shape(1, Inches(x + 0.3), Inches(2.95), Inches(1.7), Inches(0.04))
        line2.fill.solid()
        line2.fill.fore_color.rgb = hex_color(border)
        line2.line.fill.background()
        add_text_box(sl, desc, x + 0.1, 3.05, 2.1, 2.7, font_size=10, color_hex="D8B4FE", align=PP_ALIGN.CENTER)

    # ===================== СЛАЙД 8: ЗАКЛЮЧЕНИЕ =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 0.08, PURPLE, GREEN)
    add_text_box(sl, "ИТОГИ", 0.4, 0.2, 3, 0.35, font_size=9, color_hex="A78BFA")
    add_text_box(sl, "Главные выводы", 0.4, 0.55, 8, 0.6, font_size=32, bold=True, color_hex=WHITE)

    conclusions = [
        ("🧠", "Алкоголь — клеточный яд, убивающий нейроны"),
        ("🫀", "Печень разрушается необратимо (цирроз)"),
        ("❤️", "Сердце и сосуды страдают от тромбов"),
        ("🧬", "Поражение половых клеток влияет на здоровье будущих детей"),
    ]
    for i, (icon, text) in enumerate(conclusions):
        add_card(sl, 0.4, 1.4 + i * 0.95, 6.0, 0.8, CARD, BORDER)
        add_text_box(sl, icon, 0.6, 1.45 + i * 0.95, 0.6, 0.6, font_size=24)
        add_text_box(sl, text, 1.3, 1.5 + i * 0.95, 4.9, 0.55, font_size=13, bold=True, color_hex=LIGHT)

    add_card(sl, 6.8, 1.4, 6.1, 2.0, "1F0000", RED)
    add_text_box(sl, "❌", 9.5, 1.5, 0.6, 0.5, font_size=22, align=PP_ALIGN.CENTER)
    add_text_box(sl, "БЕЗВРЕДНЫХ ДОЗ НЕ СУЩЕСТВУЕТ!", 6.9, 1.95, 5.8, 0.55, font_size=16, bold=True, color_hex="FCA5A5", align=PP_ALIGN.CENTER)
    add_text_box(sl, "Особенно опасен алкоголь для подростков —\nорганизм ещё формируется",
                 6.9, 2.52, 5.8, 0.75, font_size=11, color_hex="FEE2E2", align=PP_ALIGN.CENTER)

    add_card(sl, 6.8, 3.6, 6.1, 1.7, "001A00", GREEN)
    add_text_box(sl, "📋", 9.55, 3.7, 0.6, 0.5, font_size=22, align=PP_ALIGN.CENTER)
    add_text_box(sl, "Результат работы", 6.9, 4.2, 5.8, 0.4, font_size=14, bold=True, color_hex="86EFAC", align=PP_ALIGN.CENTER)
    add_text_box(sl, "Создана памятка «Биология трезвости», которая поможет\nсверстникам сделать осознанный выбор",
                 6.9, 4.62, 5.8, 0.55, font_size=11, color_hex="D1FAE5", align=PP_ALIGN.CENTER)

    add_gradient_rect(sl, 0.4, 5.8, 12.5, 0.7, "3B0764", "831843")
    add_text_box(sl, "Знание биологии = осознанный выбор в пользу здоровья!",
                 0.5, 5.88, 12.3, 0.5, font_size=16, bold=True, color_hex=WHITE, align=PP_ALIGN.CENTER)

    # ===================== СЛАЙД 9: СПАСИБО =====================
    sl = prs.slides.add_slide(blank_layout)
    set_bg(sl, BG)
    add_gradient_rect(sl, 0, 0, 13.33, 7.5, "1A0040", "0D0020")
    add_gradient_rect(sl, 0, 0, 13.33, 7.5, "2D0060", "0D0020")

    add_text_box(sl, "🎓", 5.9, 0.5, 1.5, 1.2, font_size=54, align=PP_ALIGN.CENTER)

    add_gradient_rect(sl, 1.5, 1.75, 10.3, 1.05, "4C1D95", "831843")
    add_text_box(sl, "СПАСИБО ЗА ВНИМАНИЕ!", 1.5, 1.8, 10.3, 0.9, font_size=40, bold=True, color_hex=WHITE, align=PP_ALIGN.CENTER)

    line3 = sl.shapes.add_shape(1, Inches(4.5), Inches(2.95), Inches(4.33), Inches(0.07))
    line3.fill.gradient()
    line3.fill.gradient_angle = 0
    line3.fill.gradient_stops[0].color.rgb = hex_color(PURPLE)
    line3.fill.gradient_stops[1].color.rgb = hex_color(ORANGE)
    line3.line.fill.background()

    add_text_box(sl, "Готов ответить на ваши вопросы",
                 2.5, 3.2, 8.3, 0.5, font_size=16, color_hex="C4B5FD", align=PP_ALIGN.CENTER)

    add_card(sl, 3.7, 4.0, 6.0, 1.8, CARD, BORDER)
    add_text_box(sl, "Выполнил: Садомов Илья, 9Д класс",
                 3.9, 4.1, 5.6, 0.4, font_size=13, bold=True, color_hex="F9A8D4", align=PP_ALIGN.CENTER)
    add_text_box(sl, "Руководитель: Садомова Екатерина Геннадьевна",
                 3.9, 4.55, 5.6, 0.35, font_size=11, color_hex="D8B4FE", align=PP_ALIGN.CENTER)
    add_text_box(sl, "г. Красноярск, 2026",
                 3.9, 4.95, 5.6, 0.35, font_size=11, color_hex="A78BFA", align=PP_ALIGN.CENTER)

    # Сохраняем в буфер и кодируем base64
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode('utf-8')

    return {
        'statusCode': 200,
        'headers': {
            'Access-Control-Allow-Origin': '*',
            'Content-Type': 'application/json',
        },
        'body': f'{{"file": "{b64}", "filename": "Влияние_алкоголя_на_организм.pptx"}}'
    }
